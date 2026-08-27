import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from PIL import Image

# Core Constants
INPUT_IMAGE_PATH = "/home/pedro/Documentos/pesquisa/midias/noite_estrelada_original.jpg"

# Color Definitions
DARK_BLACK = RGBColor(0, 0, 0)
LIGHT_WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(212, 175, 55)        # #D4AF37 (Dark theme title and bullets)
LIGHT_BLUE = RGBColor(147, 197, 253) # #93C5FD (Dark theme links)
DARK_BLUE = RGBColor(15, 23, 42)     # #0F172A (Light theme title)
LIGHT_GREY_LINE = RGBColor(220, 220, 220)
DARK_GREY_LINE = RGBColor(80, 80, 80)
OFF_WHITE = RGBColor(224, 230, 237)  # #E0E6ED
SLATE_BLACK = RGBColor(30, 41, 59)   # #1E293B (Light theme text)
IFF_GREEN = RGBColor(0, 146, 63)     # #00923F (Institutional IFF Green)
CARD_BG = RGBColor(15, 15, 15)       # Very dark card background for contrast

# 1. Background Generator
def prepare_background_images():
    img = Image.open(INPUT_IMAGE_PATH).convert("RGB")
    w, h = img.size
    
    # 16:9 aspect crop
    target_ratio = 16.0 / 9.0
    current_ratio = float(w) / float(h)
    
    if current_ratio > target_ratio:
        new_w = int(h * target_ratio)
        left = (w - new_w) // 2
        right = left + new_w
        top = 0
        bottom = h
    else:
        new_h = int(w / target_ratio)
        top = (h - new_h) // 2
        bottom = top + new_h
        left = 0
        right = w
        
    cropped_img = img.crop((left, top, right, bottom)).resize((1920, 1080), Image.Resampling.LANCZOS)
    
    # Save original cropped
    cropped_img.save("/home/pedro/Downloads/_cosmic_assets/cosmic_bg_original_cropped.jpg", "JPEG", quality=85)
    
    # Save blended with black (8% for content slides)
    black_bg = Image.new("RGB", (1920, 1080), (0, 0, 0))
    faint_dark = Image.blend(black_bg, cropped_img, 0.08)
    faint_dark.save("/home/pedro/Downloads/_cosmic_assets/cosmic_bg_atual_dark.jpg", "JPEG", quality=85)
    
    # Save blended with black (35% for cover and Q&A slides - higher opacity)
    vivid_dark = Image.blend(black_bg, cropped_img, 0.35)
    vivid_dark.save("/home/pedro/Downloads/_cosmic_assets/cosmic_bg_vivid_dark.jpg", "JPEG", quality=85)
    
    # Save blended with white (8% for content slides)
    white_bg = Image.new("RGB", (1920, 1080), (255, 255, 255))
    faint_light = Image.blend(white_bg, cropped_img, 0.08)
    faint_light.save("/home/pedro/Downloads/_cosmic_assets/cosmic_bg_atual_light.jpg", "JPEG", quality=85)
    
    # Save blended with white (35% for cover and Q&A slides - higher opacity)
    vivid_light = Image.blend(white_bg, cropped_img, 0.35)
    vivid_light.save("/home/pedro/Downloads/_cosmic_assets/cosmic_bg_vivid_light.jpg", "JPEG", quality=85)
    print("Backgrounds prepared successfully!")

# Slide Transition helper
def add_transition_fade(slide):
    transition_xml = f'<p:transition {nsdecls("p")} spd="med"><p:fade/></p:transition>'
    transition_elm = parse_xml(transition_xml)
    sld = slide.element
    
    insert_before = None
    for child in sld:
        if child.tag in (f'{{{sld.nsmap["p"]}}}timing', f'{{{sld.nsmap["p"]}}}extLst'):
            insert_before = child
            break
            
    if insert_before is not None:
        insert_before.addprevious(transition_elm)
    else:
        sld.append(transition_elm)

# Element-level click Fade In animation using mainSeq
def animate_shapes(slide, shapes):
    sld = slide.element
    
    # Clear any existing timing node
    for child in list(sld):
        if child.tag == f'{{{sld.nsmap["p"]}}}timing':
            sld.remove(child)
            
    if not shapes:
        return
        
    child_xmls = []
    for i, shape in enumerate(shapes):
        seq_id = i + 1
        child_xmls.append(f"""
            <p:par>
                <p:cTn id="{seq_id * 10}" fill="hold" nodeType="clickEffect">
                    <p:stCondLst>
                        <p:cond evt="onBegin" delay="0"/>
                    </p:stCondLst>
                    <p:childTnLst>
                        <p:anim dur="500" calcmode="linear" valueType="num">
                            <p:cBhvr>
                                <p:cTn id="{seq_id * 10 + 1}" dur="500"/>
                                <p:tgtEl>
                                    <p:spTgt spid="{shape.shape_id}"/>
                                </p:tgtEl>
                                <p:attrNameLst>
                                    <p:attrName>style.opacity</p:attrName>
                                </p:attrNameLst>
                            </p:cBhvr>
                            <p:tavLst>
                                <p:tav tm="0">
                                    <p:val>
                                        <p:numVal val="0"/>
                                    </p:val>
                                </p:tav>
                                <p:tav tm="100000">
                                    <p:val>
                                        <p:numVal val="1"/>
                                    </p:val>
                                </p:tav>
                            </p:tavLst>
                        </p:anim>
                        <p:set>
                            <p:cBhvr>
                                <p:cTn id="{seq_id * 10 + 2}" dur="1" fill="hold"/>
                                <p:tgtEl>
                                    <p:spTgt spid="{shape.shape_id}"/>
                                </p:tgtEl>
                                <p:attrNameLst>
                                    <p:attrName>style.visibility</p:attrName>
                                </p:attrNameLst>
                            </p:cBhvr>
                            <p:to>
                                <p:strVal val="visible"/>
                            </p:to>
                        </p:set>
                    </p:childTnLst>
                </p:cTn>
            </p:par>
        """)
        
    children_str = "".join(child_xmls)
    
    timing_xml = f"""<p:timing {nsdecls("p")}>
        <p:tnLst>
            <p:par>
                <p:cTn id="0" dur="indefinite" nodeType="tmRoot">
                    <p:childTnLst>
                        <p:seq concurrent="1" nextAc="nextTemplate" type="sldAnimMeth">
                            <p:cTn id="1" dur="indefinite" nodeType="mainSeq">
                                <p:childTnLst>
                                    {children_str}
                                </p:childTnLst>
                            </p:cTn>
                        </p:seq>
                    </p:childTnLst>
                </p:cTn>
            </p:par>
        </p:tnLst>
    </p:timing>"""
    
    timing_elm = parse_xml(timing_xml)
    sld.append(timing_elm)

# Scaling Helper Class
class LayoutEngine:
    def __init__(self, w, h):
        self.w = w
        self.h = h
        self.scale_x = w / 13.333
        
    def get_x(self, x_in_16_9):
        return Inches(x_in_16_9 * self.scale_x)
        
    def get_y(self, y_in_16_9):
        return Inches(y_in_16_9) # Keep vertical height fixed at 7.5"
        
    def get_width(self, w_in_16_9):
        return Inches(w_in_16_9 * self.scale_x)
        
    def get_height(self, h_in_16_9):
        return Inches(h_in_16_9)

def fit_picture_aspect_ratio(slide, image_path, slide_width, slide_height):
    img = Image.open(image_path)
    img_w, img_h = img.size
    img_ratio = img_w / img_h
    slide_ratio = slide_width / slide_height
    
    if img_ratio > slide_ratio:
        # Scale to width
        new_w = slide_width
        new_h = slide_width / img_ratio
        left = Inches(0)
        top = (slide_height - new_h) / 2
    else:
        # Scale to height
        new_h = slide_height
        new_w = slide_height * img_ratio
        top = Inches(0)
        left = (slide_width - new_w) / 2
        
    slide.shapes.add_picture(image_path, left, top, new_w, new_h)

class SlideBuilder:
    def __init__(self, theme, layout_engine, total_pages=13):
        self.theme = theme # 'Preto' or 'Branco'
        self.le = layout_engine
        self.total_pages = total_pages
        
        # Colors based on theme
        if theme == 'Preto': # MWBR Preto
            self.bg_color = DARK_BLACK
            self.title_color = GOLD
            self.body_color = OFF_WHITE
            self.accent_color = GOLD
            self.link_color = GOLD
            self.card_bg = CARD_BG
            self.line_color = DARK_GREY_LINE
            self.bg_image_faint = "/home/pedro/Downloads/_cosmic_assets/cosmic_bg_atual_dark.jpg"
            self.bg_image_vivid = "/home/pedro/Downloads/_cosmic_assets/cosmic_bg_vivid_dark.jpg"
            self.qr_code = "/home/pedro/Downloads/_cosmic_assets/cosmic_qrcode_transparente_branco.png"
            self.iff_logo = "/home/pedro/Downloads/_cosmic_assets/iff/iff_bji_dark.png"
            self.header_text = "Milky Way Brazil 2026"
            self.footer_title = "Reconciliação Química Estelar • MWBR"
        elif theme == 'Institucional_IFF_Branco':
            self.bg_color = LIGHT_WHITE
            self.title_color = IFF_GREEN
            self.body_color = SLATE_BLACK
            self.accent_color = IFF_GREEN
            self.link_color = IFF_GREEN
            self.card_bg = LIGHT_WHITE
            self.line_color = LIGHT_GREY_LINE
            self.bg_image_faint = None
            self.bg_image_vivid = None
            self.qr_code = "/home/pedro/Downloads/_cosmic_assets/cosmic_qrcode_transparente_preto.png"
            self.iff_logo = "/home/pedro/Downloads/_cosmic_assets/iff/iff_bji_light.png"
            self.header_text = "Astrofísica Geral"
            self.footer_title = "IFF Bom Jesus do Itabapoana"
        elif theme == 'Institucional_IFF_Preto':
            self.bg_color = DARK_BLACK
            self.title_color = IFF_GREEN
            self.body_color = OFF_WHITE
            self.accent_color = IFF_GREEN
            self.link_color = IFF_GREEN
            self.card_bg = CARD_BG
            self.line_color = DARK_GREY_LINE
            self.bg_image_faint = None
            self.bg_image_vivid = None
            self.qr_code = "/home/pedro/Downloads/_cosmic_assets/cosmic_qrcode_transparente_branco.png"
            self.iff_logo = "/home/pedro/Downloads/_cosmic_assets/iff/iff_bji_dark.png"
            self.header_text = "Astrofísica Geral"
            self.footer_title = "IFF Bom Jesus do Itabapoana"
        else: # Branco / MWBR
            self.bg_color = LIGHT_WHITE
            self.title_color = DARK_BLUE
            self.body_color = SLATE_BLACK
            self.accent_color = IFF_GREEN
            self.link_color = IFF_GREEN
            self.card_bg = LIGHT_WHITE
            self.line_color = LIGHT_GREY_LINE
            self.bg_image_faint = "/home/pedro/Downloads/_cosmic_assets/cosmic_bg_atual_light.jpg"
            self.bg_image_vivid = "/home/pedro/Downloads/_cosmic_assets/cosmic_bg_vivid_light.jpg"
            self.qr_code = "/home/pedro/Downloads/_cosmic_assets/cosmic_qrcode_transparente_preto.png"
            self.iff_logo = "/home/pedro/Downloads/_cosmic_assets/iff/iff_bji_light.png"
            self.header_text = "Milky Way Brazil 2026"
            self.footer_title = "Reconciliação Química Estelar • MWBR"
            
        self.bg_image_original = "/home/pedro/Downloads/_cosmic_assets/cosmic_bg_original_cropped.jpg"

    def fix_theme_hyperlink_colors(self, prs):
        # Iterate package parts to locate theme xml safely
        if self.theme == 'Preto':
            hex_color = "D4AF37"
        elif self.theme == 'Institucional_IFF':
            hex_color = "00923F"
        else:
            hex_color = "0F766E"
        for part in prs.part.package.iter_parts():
            if "theme" in str(part.partname):
                theme_xml = part.blob.decode('utf-8')
                new_xml = theme_xml.replace('<a:hlink><a:srgbClr val="0000FF"/></a:hlink>', f'<a:hlink><a:srgbClr val="{hex_color}"/></a:hlink>')
                new_xml = new_xml.replace('<a:folHlink><a:srgbClr val="800080"/></a:folHlink>', f'<a:folHlink><a:srgbClr val="{hex_color}"/></a:folHlink>')
                part._blob = new_xml.encode('utf-8')

    def create_slide_base(self, prs, bg_image=None):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
        
        if bg_image and os.path.exists(bg_image):
            slide.shapes.add_picture(bg_image, 0, 0, prs.slide_width, prs.slide_height)
            
        return slide

    def add_slide_headers(self, slide, title_text, prs):
        # Header line: spans from absolute left edge (0) to absolute right edge of slide
        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.65), prs.slide_width, Inches(0.01))
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.line_color
        divider.line.fill.background()

        # 1. Top left official IFF Bom Jesus do Itabapoana Logo Image
        if os.path.exists(self.iff_logo):
            # Pos=(0.8", 0.15") Height=0.45" Width proportional (~1.2in)
            slide.shapes.add_picture(self.iff_logo, self.le.get_x(0.8), Inches(0.12), height=Inches(0.48))
        else:
            logo_box = slide.shapes.add_textbox(self.le.get_x(0.8), self.le.get_y(0.2), self.le.get_width(4.0), self.le.get_height(0.35))
            tf_logo = logo_box.text_frame
            tf_logo.margin_left = tf_logo.margin_top = tf_logo.margin_right = tf_logo.margin_bottom = 0
            p_logo = tf_logo.paragraphs[0]
            p_logo.text = "IFF Bom Jesus do Itabapoana"
            p_logo.font.name = "Arial"
            p_logo.font.size = Pt(14)
            p_logo.font.bold = True
            p_logo.font.color.rgb = self.accent_color
        
        # 2. Top right event header (larger font size: 12pt, aligned to the right margin)
        right_x = prs.slide_width - Inches(0.8) - self.le.get_width(5.0)
        hdr_box = slide.shapes.add_textbox(right_x, self.le.get_y(0.2), self.le.get_width(5.0), self.le.get_height(0.35))
        tf_hdr = hdr_box.text_frame
        tf_hdr.margin_left = tf_hdr.margin_top = tf_hdr.margin_right = tf_hdr.margin_bottom = 0
        p_hdr = tf_hdr.paragraphs[0]
        p_hdr.alignment = PP_ALIGN.RIGHT
        p_hdr.text = self.header_text
        p_hdr.font.name = "Arial"
        p_hdr.font.size = Pt(12)
        p_hdr.font.color.rgb = self.body_color
        
        # 3. Slide Title
        title_box = slide.shapes.add_textbox(self.le.get_x(0.8), self.le.get_y(0.85), self.le.get_width(11.733), self.le.get_height(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.title_color

    def add_slide_footer(self, slide, prs, page_num=None):
        # Footer Divider Line: spans from absolute left edge (0) to absolute right edge of slide, closer to bottom (7.20")
        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.20), prs.slide_width, Inches(0.01))
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.line_color
        divider.line.fill.background()
        
        # Left footer (Presenter + hyperlink)
        left_box = slide.shapes.add_textbox(self.le.get_x(0.8), self.le.get_y(7.24), self.le.get_width(4.0), self.le.get_height(0.2))
        tf_l = left_box.text_frame
        tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
        p_l = tf_l.paragraphs[0]
        run_l = p_l.add_run()
        run_l.text = "Pedro H. R. de Andrade"
        run_l.hyperlink.address = "mailto:pedroiff0@gmail.com"
        run_l.font.name = "Arial"
        run_l.font.size = Pt(9)
        run_l.font.color.rgb = self.body_color
        
        # Center footer (Short Title / Institution)
        center_x = (prs.slide_width - self.le.get_width(5.0)) / 2
        center_box = slide.shapes.add_textbox(center_x, self.le.get_y(7.24), self.le.get_width(5.0), self.le.get_height(0.2))
        tf_c = center_box.text_frame
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        run_c = p_c.add_run()
        run_c.text = self.footer_title
        run_c.font.name = "Arial"
        run_c.font.size = Pt(9)
        run_c.font.color.rgb = self.body_color
        
        # Right footer (AUTOMATIC Slide Number Field + Total Pages)
        # In Keynote/PowerPoint, slide number placeholders MUST have <p:ph type="sldNum"/> in cNvPr/nvPr!
        right_x = prs.slide_width - Inches(0.8) - self.le.get_width(4.0)
        right_box = slide.shapes.add_textbox(right_x, self.le.get_y(7.24), self.le.get_width(4.0), self.le.get_height(0.2))
        
        # Set placeholder type sldNum in XML so Keynote recognizes it as the official Slide Number object
        try:
            nvPr = right_box._element.nvSpPr.nvPr
            ph_xml = f'<p:ph {nsdecls("p")} type="sldNum" sz="quarter"/>'
            nvPr.append(parse_xml(ph_xml))
        except Exception:
            pass

        tf_r = right_box.text_frame
        tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
        p_r = tf_r.paragraphs[0]
        p_r.alignment = PP_ALIGN.RIGHT
        
        hex_color = self.theme_hex_color()
        page_str = str(page_num) if page_num is not None else "‹#›"
        
        fld_xml = f'''<a:fld {nsdecls("a")} id="{{53AF88E2-2E5C-47E5-9D67-7D15B9A608F7}}" type="slidenum">
            <a:rPr sz="900" face="Arial">
                <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
            </a:rPr>
            <a:t>{page_str}</a:t>
        </a:fld>'''
        fld_elm = parse_xml(fld_xml)
        p_r._p.append(fld_elm)
        
        run_total_xml = f'''<a:r {nsdecls("a")}>
            <a:rPr sz="900" face="Arial">
                <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
            </a:rPr>
            <a:t> / {self.total_pages}</a:t>
        </a:r>'''
        run_total_elm = parse_xml(run_total_xml)
        p_r._p.append(run_total_elm)

    def theme_hex_color(self):
        c = self.body_color
        return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

    def add_bullet_points(self, tf, bullets):
        for bold_text, normal_text in bullets:
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.space_after = Pt(12)
            
            run_bold = p.add_run()
            run_bold.text = bold_text
            run_bold.font.name = "Arial"
            run_bold.font.size = Pt(14)
            run_bold.font.bold = True
            run_bold.font.color.rgb = self.accent_color
            
            run_norm = p.add_run()
            run_norm.text = normal_text
            run_norm.font.name = "Arial"
            run_norm.font.size = Pt(14)
            run_norm.font.color.rgb = self.body_color

    def add_animated_bullets(self, slide, bullets, x, y, width, line_gap, height=0.9):
        shapes = []
        curr_y = y
        for bold_text, normal_text in bullets:
            box = slide.shapes.add_textbox(self.le.get_x(x), self.le.get_y(curr_y), self.le.get_width(width), self.le.get_height(height))
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            
            p = tf.paragraphs[0]
            p.space_after = Pt(0)
            
            run_bold = p.add_run()
            run_bold.text = bold_text
            run_bold.font.name = "Arial"
            run_bold.font.size = Pt(14)
            run_bold.font.bold = True
            run_bold.font.color.rgb = self.accent_color
            
            run_norm = p.add_run()
            run_norm.text = normal_text
            run_norm.font.name = "Arial"
            run_norm.font.size = Pt(14)
            run_norm.font.color.rgb = self.body_color
            
            shapes.append(box)
            curr_y += line_gap
        return shapes

    def add_animated_references(self, slide, refs, x, y, width, line_gap, height=0.6):
        shapes = []
        curr_y = y
        for bold_text, normal_text in refs:
            box = slide.shapes.add_textbox(self.le.get_x(x), self.le.get_y(curr_y), self.le.get_width(width), self.le.get_height(height))
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            
            p = tf.paragraphs[0]
            p.space_after = Pt(0)
            
            run_bold = p.add_run()
            run_bold.text = bold_text
            run_bold.font.name = "Arial"
            run_bold.font.size = Pt(13)
            run_bold.font.bold = True
            run_bold.font.color.rgb = self.accent_color
            
            run_norm = p.add_run()
            run_norm.text = normal_text
            run_norm.font.name = "Arial"
            run_norm.font.size = Pt(13)
            run_norm.font.color.rgb = self.body_color
            
            shapes.append(box)
            curr_y += line_gap
        return shapes

    def add_citation(self, slide, citation_text):
        cite_box = slide.shapes.add_textbox(self.le.get_x(0.8), self.le.get_y(6.6), self.le.get_width(11.733), self.le.get_height(0.3))
        tf = cite_box.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = f"Ref: {citation_text}"
        p.font.name = "Arial"
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = self.body_color

    # Slide 1: Capa (No footer / No pagination shown!)
    def build_title_slide(self, prs, presenter_name, presenter_email):
        slide = self.create_slide_base(prs, self.bg_image_vivid)
        
        if os.path.exists(self.iff_logo):
            slide.shapes.add_picture(self.iff_logo, self.le.get_x(0.8), Inches(0.5), height=Inches(0.95))
        else:
            logo_box = slide.shapes.add_textbox(self.le.get_x(0.8), self.le.get_y(0.5), self.le.get_width(4.5), self.le.get_height(0.6))
            tf_logo = logo_box.text_frame
            p_logo = tf_logo.paragraphs[0]
            p_logo.text = "IFF Bom Jesus do Itabapoana"
            p_logo.font.name = "Arial"
            p_logo.font.size = Pt(18)
            p_logo.font.bold = True
            p_logo.font.color.rgb = self.accent_color
        
        title_box = slide.shapes.add_textbox(self.le.get_x(0.8), self.le.get_y(1.8), self.le.get_width(11.733), self.le.get_height(2.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "Reconciliação Química no Disco Galáctico"
        p.font.name = "Arial"
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = self.title_color
        
        p_sub = tf.add_paragraph()
        p_sub.text = "Abundâncias Estelares de Anãs M e Estrelas FGK face à Evolução Química Galáctica (GCE)"
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(18)
        p_sub.font.color.rgb = self.body_color
        p_sub.space_before = Pt(14)
        
        info_box = slide.shapes.add_textbox(self.le.get_x(0.8), self.le.get_y(4.8), self.le.get_width(11.733), self.le.get_height(1.8))
        tf_info = info_box.text_frame
        tf_info.word_wrap = True
        
        p_pres = tf_info.paragraphs[0]
        p_pres.text = f"Apresentador: {presenter_name}"
        p_pres.font.name = "Arial"
        p_pres.font.size = Pt(16)
        p_pres.font.bold = True
        p_pres.font.color.rgb = self.body_color
        
        p_email = tf_info.add_paragraph()
        p_email.text = f"Contato: {presenter_email}"
        p_email.font.name = "Arial"
        p_email.font.size = Pt(14)
        p_email.font.color.rgb = self.link_color
        
        p_art = tf_info.add_paragraph()
        p_art.text = "Artigo: Reconciliation of M Dwarf and FGK Stellar Abundances with Galactic Chemical Evolution (ApJ, 2026)\nAutores: T. C. L. Trueman, H. S. Wang, S. J. Mojzsis, D. Turrini, K. Heng, M. Pignatari"
        p_art.font.name = "Arial"
        p_art.font.size = Pt(12)
        p_art.font.color.rgb = self.body_color
        p_art.space_before = Pt(10)
        
        # No footer call here! Cover slide is clean.
        add_transition_fade(slide)

    # Slide 2: Sumário
    def build_summary_slide(self, prs, page_num=1):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "Sumário da Apresentação", prs)
        
        bullets = [
            ("1. Introdução e Contexto: ", "O conflito histórico de abundâncias químicas no disco local."),
            ("2. Desafios Espectroscópicos: ", "A complexidade das atmosferas frias e bandas de TiO."),
            ("3. Metodologia Observacional e Teórica: ", "Estudos benchmark 3D NLTE e modelos GCE (OMEGA+)."),
            ("4. Reconciliação Química: ", "Resultados de C, O, Na e Ni e o enigma persistente do Titânio."),
            ("5. Conexão Exoplanetária e Conclusões: ", "Impacto na geofísica de planetas rochosos e a missão Ariel."),
            ("6. Referências Bibliográficas: ", "Literatura fundamental utilizada no trabalho.")
        ]
        shapes = self.add_animated_bullets(slide, bullets, x=0.8, y=1.9, width=11.733, line_gap=0.75, height=0.6)
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        # Entrance animation for the summary list
        animate_shapes(slide, shapes)

    # Slide 3: Introdução e Contexto
    def build_intro_slide(self, prs, page_num=2):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "1. Introdução e Contexto", prs)
        
        bullets = [
            ("Anãs M na Galáxia: ", "Constituem cerca de 70% da população estelar da Via Láctea, alvos principais para buscas de exoplanetas."),
            ("Preservação Química: ", "Sua evolução ultra lenta preserva na fotosfera a assinatura intacta da nuvem molecular original."),
            ("Problema Científico: ", "Estudos anteriores encontravam discrepâncias químicas severas entre anãs M e FGK no disco local.")
        ]
        
        # Left Text (Separate textbox shapes for paragraph transitions)
        left_shapes = self.add_animated_bullets(slide, bullets, x=0.8, y=1.9, width=5.6, line_gap=1.3, height=1.1)
        
        # Right Card
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.le.get_x(6.8), self.le.get_y(1.9), self.le.get_width(5.7), self.le.get_height(4.3))
        card.fill.solid()
        card.fill.fore_color.rgb = self.card_bg
        card.line.color.rgb = self.accent_color
        card.line.width = Pt(1.0)
        
        tf_r = card.text_frame
        tf_r.word_wrap = True
        tf_r.margin_left = tf_r.margin_right = Inches(0.3)
        tf_r.margin_top = tf_r.margin_bottom = Inches(0.3)
        
        p = tf_r.paragraphs[0]
        p.text = "A Controvérsia Química"
        p.font.name = "Arial"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.accent_color
        p.space_after = Pt(14)
        
        card_bullets = [
            ("Modelagem Física: ", "Vieses sistemáticos de 1D LTE distorciam a análise espectral nas anãs M."),
            ("Questão-Chave: ", "Isso indica enriquecimento anômalo ou falhas de espectroscopia?")
        ]
        
        for bold_text, normal_text in card_bullets:
            p_c = tf_r.add_paragraph()
            p_c.space_after = Pt(12)
            rb = p_c.add_run()
            rb.text = bold_text
            rb.font.name = "Arial"
            rb.font.bold = True
            rb.font.size = Pt(14)
            rb.font.color.rgb = self.accent_color
            
            rn = p_c.add_run()
            rn.text = normal_text
            rn.font.name = "Arial"
            rn.font.size = Pt(14)
            rn.font.color.rgb = self.body_color
            
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        # Click animations for left text bullets first, then right card
        animate_shapes(slide, left_shapes + [card])

    # Slide 4: Desafios Espectroscópicos
    def build_challenges_slide(self, prs, page_num=3):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "2. Desafios Espectroscópicos em Estrelas Frias", prs)
        
        bullets = [
            ("Temperaturas Baixas (< 4000 K): ", "Permitem condensação molecular na fotosfera, saturando o espectro."),
            ("Bandas de TiO: ", "Óxido de Titânio polui e bloqueia o contínuo óptico, cobrindo linhas atômicas livres."),
            ("Vieses de 1D LTE: ", "Suposição clássica de LTE unidimensional distorce abundâncias obtidas de linhas de ressonância."),
            ("Importância de NLTE 3D: ", "Modelos avançados são cruciais para isolar metalicidades e elementos refratários como Na e Al.")
        ]
        shapes = self.add_animated_bullets(slide, bullets, x=0.8, y=1.9, width=11.733, line_gap=1.05, height=0.9)
        self.add_citation(slide, "Amarsi et al. (2019), Bergemann et al. (2017)")
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        # Animate each bullet sequentially
        animate_shapes(slide, shapes)

    # Slide 5: Metodologia e Amostra (Text + Image placeholder)
    def build_observational_slide(self, prs, page_num=4):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "3. Metodologia Observacional: Amostra", prs)
        
        bullets = [
            ("Amostra Benchmark: ", "47 estrelas anãs M selecionadas dos catálogos APOGEE, CARMENES e Subaru/IRD."),
            ("Homogeneidade: ", "Uso exclusivo de atmosferas modelo MARCS para calibração espectral."),
            ("FGK de Referência: ", "Dados compilados e corrigidos para efeitos tridimensionais Não-LTE no disco local."),
            ("Faixa de Metalicidade: ", "Foco no disco fino galáctico cobrindo -0.5 <= [Fe/H] <= +0.4.")
        ]
        left_shapes = self.add_animated_bullets(slide, bullets, x=0.8, y=1.9, width=5.8, line_gap=1.05, height=0.9)
        self.add_citation(slide, "Gustafsson et al. (2008), Trueman et al. (2026)")
        
        # Right Image Layout Placeholder
        img_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.le.get_x(7.0), self.le.get_y(1.9), self.le.get_width(5.5), self.le.get_height(4.3))
        img_card.fill.solid()
        img_card.fill.fore_color.rgb = self.card_bg
        img_card.line.color.rgb = self.accent_color
        img_card.line.width = Pt(1.0)
        
        tf_img = img_card.text_frame
        tf_img.word_wrap = True
        p_img = tf_img.paragraphs[0]
        p_img.alignment = PP_ALIGN.CENTER
        p_img.text = "\n\n\n\n[ Espaço Reservado para Figura: Espectros / Diagrama HR ]"
        p_img.font.name = "Arial"
        p_img.font.size = Pt(14)
        p_img.font.bold = True
        p_img.font.color.rgb = self.accent_color
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        # Click animates text column first, then image layout card
        animate_shapes(slide, left_shapes + [img_card])

    # Slide 6: Modelagem Teórica GCE
    def build_gce_slide(self, prs, page_num=5):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "4. Metodologia Teórica: Modelos GCE", prs)
        
        bullets = [
            ("Código OMEGA+: ", "Simulações de enriquecimento químico baseadas em modelo de duas zonas (galáxia e meio circumgaláctico)."),
            ("Rendimentos de Nucleossíntese: ", "Modelos comparativos de supernovas core-collapse (CCSN) de Nomoto (2013) e Limongi & Chieffi (2018)."),
            ("Supernovas Tipo Ia: ", "Delay-Time Distribution (DTD) calibrado com platô de 1 Gyr para ajuste dos elementos do grupo do ferro."),
            ("Dispersão no Disco: ", "Simulação de 6 trilhas independentes mapeando o espalhamento de metalicidade solar local.")
        ]
        shapes = self.add_animated_bullets(slide, bullets, x=0.8, y=1.9, width=11.733, line_gap=1.05, height=0.9)
        self.add_citation(slide, "Nomoto et al. (2013), Limongi & Chieffi (2018)")
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        animate_shapes(slide, shapes)

    # Slide 7: Resultados Reconciliação (Text + Image placeholder)
    def build_results_slide(self, prs, page_num=6):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "5. Resultados: Reconciliação Química", prs)
        
        bullets = [
            ("Slopes Equivalentes: ", "Abundâncias de [C/H], [O/H], [Na/Fe] e [Ni/Fe] nas anãs M e FGK apresentam a mesma inclinação."),
            ("Ajuste por Regressão ODR: ", "Sobreposição completa de 1-sigma nas bandas de confiança estatística."),
            ("Resolução de Divergências: ", "Estudos anteriores divergiam porque não corrigiam os dados de FGK com modelos 3D NLTE."),
            ("Sucesso do OMEGA+: ", "Excelente acordo entre previsões dos modelos GCE e observações estelares.")
        ]
        left_shapes = self.add_animated_bullets(slide, bullets, x=0.8, y=1.9, width=5.8, line_gap=1.05, height=0.9)
        
        # Right Image
        img_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.le.get_x(7.0), self.le.get_y(1.9), self.le.get_width(5.5), self.le.get_height(4.3))
        img_card.fill.solid()
        img_card.fill.fore_color.rgb = self.card_bg
        img_card.line.color.rgb = self.accent_color
        img_card.line.width = Pt(1.0)
        
        tf_img = img_card.text_frame
        tf_img.word_wrap = True
        p_img = tf_img.paragraphs[0]
        p_img.alignment = PP_ALIGN.CENTER
        p_img.text = "\n\n\n\n[ Espaço Reservado para Figura: Regressão ODR / Abundâncias ]"
        p_img.font.name = "Arial"
        p_img.font.size = Pt(14)
        p_img.font.bold = True
        p_img.font.color.rgb = self.accent_color
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        animate_shapes(slide, left_shapes + [img_card])

    # Slide 8: A Exceção do Titânio
    def build_titanium_slide(self, prs, page_num=7):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "6. A Exceção Singular do Titânio ([Ti/H])", prs)
        
        bullets = [
            ("Incompatibilidade Estatística (p ~ 0.00): ", "O slope obtido para [Ti/H] em anãs M diverge severamente em relação às estrelas FGK."),
            ("Gargalo de Condensação: ", "Em atmosferas frias de anãs M, o Ti é incorporado na formação de TiO, reduzindo as linhas atômicas livres."),
            ("Vieses em Espectros: ", "Uso de linhas de Ti I atômico sem devidas correções de Não-LTE causa subestimação sistemática de abundâncias."),
            ("Causa Identificada: ", "A discrepância é de natureza metodológica observacional, e não reflete uma anomalia química real galáctica.")
        ]
        shapes = self.add_animated_bullets(slide, bullets, x=0.8, y=1.9, width=11.733, line_gap=1.05, height=0.9)
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        animate_shapes(slide, shapes)

    # Slide 9: Conexão com Exoplanetas
    def build_planets_slide(self, prs, page_num=8):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "7. Implicações na Caracterização Exoplanetária", prs)
        
        bullets = [
            ("Composição do Disco: ", "As abundâncias estelares regulam as frações de refratários e voláteis no disco protoestelar."),
            ("Geofísica de Mundos Rochosos: ", "Razões chaves como [Mg/Si] e [C/O] determinam mineralogia do manto e tamanho de núcleos de ferro planetários."),
            ("Parâmetro Ariel: ", "A missão espacial Ariel usará a composição química estelar para estimar a geodinâmica e a atmosfera dos exoplanetas."),
            ("GCE como Benchmark: ", "A consistência valida o uso de modelos galácticos para normalizar composições em surveys exoplanetários amplos.")
        ]
        shapes = self.add_animated_bullets(slide, bullets, x=0.8, y=1.9, width=11.733, line_gap=1.05, height=0.9)
        self.add_citation(slide, "Ariel White Paper (2021)")
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        animate_shapes(slide, shapes)

    # Slide 10: Conclusões e Perspectivas Futuras
    def build_conclusions_slide(self, prs, page_num=9):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "8. Conclusões e Perspectivas Futuras", prs)
        
        bullets = [
            ("Coerência Química: ", "Anãs M e estrelas FGK na vizinhança solar compartilham da mesma história de enriquecimento químico galáctico."),
            ("Superação de Vieses: ", "Modelos tridimensionais fora de LTE são cruciais para reestabelecer concordâncias espectroscópicas."),
            ("Calibração Teórica: ", "Códigos de evolução química como o OMEGA+ servem como verificação independente de discrepâncias espectrais."),
            ("Apoio a Exoplanetas: ", "A determinação refinada das abundâncias estelares garante inferências precisas sobre interiores de exoplanetas.")
        ]
        shapes = self.add_animated_bullets(slide, bullets, x=0.8, y=1.9, width=11.733, line_gap=1.05, height=0.9)
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        animate_shapes(slide, shapes)

    # Slide 11: Síntese do Conteúdo
    def build_synthesis_slide(self, prs, page_num=10):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "9. Síntese do Conteúdo", prs)
        
        bullets = [
            ("Modelagem 3D NLTE: ", "Fundamental para eliminar discrepâncias artificiais nas anãs M."),
            ("Harmonia Química Local: ", "Abundâncias de C, O, Na e Ni no disco local são homogêneas."),
            ("Desvio Metodológico: ", "O desvio de [Ti/H] é físico-químico (TiO e LTE), não nucleossintético."),
            ("Base para Exoplanetas: ", "Validação dos modelos GCE como referência para caracterizar outros mundos.")
        ]
        shapes = self.add_animated_bullets(slide, bullets, x=0.8, y=1.9, width=11.733, line_gap=1.05, height=0.9)
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        animate_shapes(slide, shapes)

    # Slide 11: Referências
    def build_references_slide(self, prs, page_num=11):
        slide = self.create_slide_base(prs, None)
        self.add_slide_headers(slide, "Referências Bibliográficas", prs)
        
        refs = [
            ("• Trueman, T. C. L., et al. 2026, ", "ApJ, 1024, L15 (Paper base do MWBR)"),
            ("• Amarsi, A. M., et al. 2019, ", "A&A, 624, A111 (Correções Não-LTE 3D para estrelas FGK)"),
            ("• Bergemann, M., et al. 2017, ", "MNRAS, 471, 331 (Modelos de atmosferas estelares em NLTE)"),
            ("• Nomoto, K., et al. 2013, ", "ARA&A, 51, 457 (Rendimentos químicos galácticos para CCSN)"),
            ("• Limongi, M., & Chieffi, A. 2018, ", "ApJS, 237, 13 (Nucleossíntese em estrelas massivas rotativas)"),
            ("• Gustafsson, B., et al. 2008, ", "A&A, 486, 951 (Grade de atmosferas modelo MARCS)")
        ]
        shapes = self.add_animated_references(slide, refs, x=0.8, y=1.9, width=11.733, line_gap=0.75, height=0.6)
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        animate_shapes(slide, shapes)

    # Slide 13: Obrigado / Q&A (Matches layout from copy.pptx)
    def build_qa_slide(self, prs, presenter_name, presenter_email, website_url, page_num=12):
        slide = self.create_slide_base(prs, self.bg_image_vivid)
        
        # 1. Header Box (Obrigado! / Dúvidas?)
        header_box = slide.shapes.add_textbox(Inches(0.73), Inches(0.65), Inches(4.1), Inches(2.0))
        tf_h = header_box.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
        
        p1 = tf_h.paragraphs[0]
        p1.text = "Obrigado!"
        p1.font.name = "Arial"
        p1.font.size = Pt(36)
        p1.font.bold = True
        p1.font.color.rgb = self.title_color
        p1.space_after = Pt(8)
        
        p2 = tf_h.add_paragraph()
        p2.text = "Dúvidas?"
        p2.font.name = "Arial"
        p2.font.size = Pt(32)
        p2.font.bold = True
        p2.font.color.rgb = LIGHT_WHITE if self.theme == 'Preto' else SLATE_BLACK
        
        # 2. Presenter Contacts Box (Left)
        contact_box = slide.shapes.add_textbox(Inches(0.70), Inches(3.12), Inches(5.4), Inches(1.3))
        tf_c = contact_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        
        pc1 = tf_c.paragraphs[0]
        pc1.text = presenter_name
        pc1.font.name = "Arial"
        pc1.font.size = Pt(22)
        pc1.font.bold = True
        pc1.font.color.rgb = LIGHT_WHITE if self.theme == 'Preto' else SLATE_BLACK
        pc1.space_after = Pt(8)
        
        pc2 = tf_c.add_paragraph()
        pc2.text = f"E-mail: {presenter_email}"
        pc2.font.name = "Arial"
        pc2.font.size = Pt(16)
        pc2.font.color.rgb = self.link_color
        
        # 3. Giant Right QR Code (Full Height from Y=0.0" to Y=7.2")
        qr_shape = None
        if os.path.exists(self.qr_code):
            # Pos=(6.16", -0.03") Size=(7.20" x 7.20")
            qr_left = Inches(6.16)
            qr_top = Inches(0.0)
            qr_size = Inches(7.20)
            qr_shape = slide.shapes.add_picture(self.qr_code, qr_left, qr_top, qr_size, qr_size)
            
        # 4. Website Link Box (Bottom Right under QR Code, above footer)
        web_box = slide.shapes.add_textbox(Inches(6.91), Inches(6.70), Inches(5.7), Inches(0.48))
        tf_w = web_box.text_frame
        tf_w.word_wrap = True
        tf_w.margin_left = tf_w.margin_top = tf_w.margin_right = tf_w.margin_bottom = 0
        pw = tf_w.paragraphs[0]
        pw.alignment = PP_ALIGN.CENTER
        
        run_w = pw.add_run()
        run_w.text = "www.phrandrade.com"
        run_w.hyperlink.address = "https://www.phrandrade.com"
        run_w.font.name = "Arial"
        run_w.font.size = Pt(14)
        run_w.font.italic = True
        run_w.font.color.rgb = self.title_color
        
        self.add_slide_footer(slide, prs, page_num=page_num)
        add_transition_fade(slide)
        
        shapes_to_animate = [header_box, contact_box, web_box]
        if qr_shape is not None:
            shapes_to_animate.append(qr_shape)
        animate_shapes(slide, shapes_to_animate)

    # Slide 13: Final Slide (Original painting fitted without cropping, no footer!)
    def build_final_slide(self, prs):
        slide = self.create_slide_base(prs, None)
        # Correctly scale the image to fit slide dimensions maintaining aspect ratio
        fit_picture_aspect_ratio(slide, self.bg_image_original, prs.slide_width, prs.slide_height)
        add_transition_fade(slide)

def main():
    prepare_background_images()
    
    aspect_ratios = [
        ('16_9', 13.333, 7.5),
        ('4_3', 10.0, 7.5),
        ('16_10', 12.0, 7.5)
    ]
    
    themes = ['Preto', 'Branco', 'Institucional_IFF_Branco', 'Institucional_IFF_Preto']
    
    for aspect_name, sw, sh in aspect_ratios:
        for theme in themes:
            prs = Presentation()
            prs.slide_width = Inches(sw)
            prs.slide_height = Inches(sh)
            # Set the starting slide number to 0 so that:
            # - Slide 1 (Capa) is page 0 (and its footer is hidden)
            # - Slide 2 (Summary) automatically numbers to 1 (showing 1 / 13)
            # - Slide 12 (Q&A) numbers to 11 (showing 11 / 13)
            # - Slide 13 (Final) is page 12 (and its footer is hidden)
            prs.element.set('firstSlideNum', '0')
            
            le = LayoutEngine(sw, sh)
            builder = SlideBuilder(theme, le, total_pages=14) # total pages is 14
            
            # Slide 1 (Cover - no footer/pagination)
            builder.build_title_slide(prs, "Pedro H. R. de Andrade", "pedroiff0@gmail.com")
            # Slide 2 (Summary - Page 1/14)
            builder.build_summary_slide(prs, page_num=1)
            # Slide 3 (Intro - Page 2/14)
            builder.build_intro_slide(prs, page_num=2)
            # Slide 4 (Challenges - Page 3/14)
            builder.build_challenges_slide(prs, page_num=3)
            # Slide 5 (Observational - Page 4/14)
            builder.build_observational_slide(prs, page_num=4)
            # Slide 6 (GCE - Page 5/14)
            builder.build_gce_slide(prs, page_num=5)
            # Slide 7 (Results - Page 6/14)
            builder.build_results_slide(prs, page_num=6)
            # Slide 8 (Titanium - Page 7/14)
            builder.build_titanium_slide(prs, page_num=7)
            # Slide 9 (Exoplanets - Page 8/14)
            builder.build_planets_slide(prs, page_num=8)
            # Slide 10 (Conclusions - Page 9/14)
            builder.build_conclusions_slide(prs, page_num=9)
            # Slide 11 (Synthesis - Page 10/14)
            builder.build_synthesis_slide(prs, page_num=10)
            # Slide 12 (References - Page 11/14)
            builder.build_references_slide(prs, page_num=11)
            # Slide 13 (Q&A - Page 12/14)
            builder.build_qa_slide(prs, "Pedro H. R. de Andrade", "pedroiff0@gmail.com", "www.phrandrade.com", page_num=12)
            # Slide 14 (Final - no footer)
            builder.build_final_slide(prs)
            
            # Fix theme XML hyperlink colors so Keynote never displays blue!
            builder.fix_theme_hyperlink_colors(prs)
            
            output_dir = "/home/pedro/Downloads/slides"
            os.makedirs(output_dir, exist_ok=True)
            output_name = os.path.join(output_dir, f"Slides_Artigo_MWBR_{aspect_name}_{theme}.pptx")
            prs.save(output_name)
            print(f"Generated: {output_name}")

if __name__ == "__main__":
    main()
